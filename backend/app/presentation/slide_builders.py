"""Individual slide type builders for the PPTX report."""
from __future__ import annotations
import re
from io import BytesIO
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.oxml.ns import qn

from .styles import *


def _set_rtl(paragraph):
    """Set RTL direction on a python-pptx paragraph via XML manipulation."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", "r")


def _set_cs_font(run):
    """Set complex script (Arabic) font on a run for proper Arabic rendering in PPTX."""
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", PPTX_FONT_CS)


def _add_text_box(
    slide: Slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 14,
    font_color=CLR_NAVY,
    bold: bool = False,
    alignment=PP_ALIGN.RIGHT,
) -> None:
    """Add an RTL Arabic text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = PPTX_FONT
    p.alignment = alignment
    _set_rtl(p)
    # Set complex script font for Arabic rendering
    for run in p.runs:
        _set_cs_font(run)


def _parse_bold_runs(text: str) -> list[dict]:
    """Split text on **bold** markers into run specs.

    Returns a list of dicts: [{"text": "...", "bold": True/False}, ...]
    """
    runs = []
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            runs.append({"text": part[2:-2], "bold": True})
        else:
            runs.append({"text": part, "bold": False})
    return runs


def _parse_markdown_to_paragraphs(text: str) -> list[dict]:
    """Parse markdown text into structured paragraph specs.

    Returns list of dicts:
      {"type": "header", "text": "...", "level": 1-3}
      {"type": "bullet", "text": "...", "runs": [...]}
      {"type": "text", "text": "...", "runs": [...]}
    """
    paragraphs = []
    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            paragraphs.append({"type": "text", "text": "", "runs": []})
            continue

        # Headers: ### Header, ## Header, # Header
        header_match = re.match(r'^(#{1,3})\s+(.+)$', stripped)
        if header_match:
            level = len(header_match.group(1))
            paragraphs.append({
                "type": "header",
                "text": header_match.group(2),
                "level": level,
                "runs": _parse_bold_runs(header_match.group(2)),
            })
            continue

        # Bullets: - item or * item or numbered: 1. item
        bullet_match = re.match(r'^[-*]\s+(.+)$', stripped)
        if not bullet_match:
            bullet_match = re.match(r'^\d+[.)]\s+(.+)$', stripped)
        if bullet_match:
            content = bullet_match.group(1)
            paragraphs.append({
                "type": "bullet",
                "text": content,
                "runs": _parse_bold_runs(content),
            })
            continue

        # Regular text with possible bold
        paragraphs.append({
            "type": "text",
            "text": stripped,
            "runs": _parse_bold_runs(stripped),
        })

    return paragraphs


def _add_rich_text_box(
    slide: Slide,
    left,
    top,
    width,
    height,
    markdown_text: str,
    base_font_size: int = 14,
    font_color=CLR_DARK_GRAY,
) -> None:
    """Build a formatted text box from markdown text with headers, bullets, bold."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs = _parse_markdown_to_paragraphs(markdown_text)

    for i, pspec in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        _set_rtl(p)

        if pspec["type"] == "header":
            size_map = {1: base_font_size + 6, 2: base_font_size + 4, 3: base_font_size + 2}
            font_size = size_map.get(pspec["level"], base_font_size + 2)
            run = p.add_run()
            run.text = pspec["text"]
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = CLR_NAVY
            run.font.name = PPTX_FONT
            _set_cs_font(run)
            p.space_before = Pt(8)

        elif pspec["type"] == "bullet":
            # Gold bullet character prefix
            bullet_run = p.add_run()
            bullet_run.text = "\u2022 "
            bullet_run.font.size = Pt(base_font_size)
            bullet_run.font.color.rgb = CLR_GOLD
            bullet_run.font.name = PPTX_FONT
            _set_cs_font(bullet_run)

            for rspec in pspec.get("runs", []):
                run = p.add_run()
                run.text = rspec["text"]
                run.font.size = Pt(base_font_size)
                run.font.color.rgb = font_color
                run.font.bold = rspec.get("bold", False)
                run.font.name = PPTX_FONT
                _set_cs_font(run)
            p.space_before = Pt(2)

        else:
            # Regular text
            for rspec in pspec.get("runs", []):
                run = p.add_run()
                run.text = rspec["text"]
                run.font.size = Pt(base_font_size)
                run.font.color.rgb = font_color
                run.font.bold = rspec.get("bold", False)
                run.font.name = PPTX_FONT
                _set_cs_font(run)


def build_title_slide(
    prs: Presentation, area_name: str, total_records: int, date_range: str
) -> Slide:
    """Slide 1: Title slide with area name, totals, date range."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CLR_NAVY

    # Gold decorative line
    line = slide.shapes.add_shape(1, Inches(1), Inches(3.4), Inches(11.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_GOLD
    line.line.fill.background()

    # Title
    _add_text_box(
        slide,
        Inches(1),
        Inches(1.5),
        Inches(11.333),
        Inches(1.2),
        "تقرير المسح العمراني",
        font_size=36,
        font_color=CLR_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Area name
    _add_text_box(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11.333),
        Inches(0.8),
        area_name,
        font_size=28,
        font_color=CLR_GOLD,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Stats line
    _add_text_box(
        slide,
        Inches(1),
        Inches(3.8),
        Inches(11.333),
        Inches(0.6),
        f"إجمالي العناوين المسحية: {total_records}  |  الفترة: {date_range}",
        font_size=16,
        font_color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Footer
    _add_text_box(
        slide,
        Inches(1),
        Inches(6.0),
        Inches(11.333),
        Inches(0.5),
        "رؤية المملكة العربية السعودية 2030",
        font_size=14,
        font_color=CLR_GOLD,
        alignment=PP_ALIGN.CENTER,
    )

    return slide


def build_summary_slide(
    prs: Presentation, stats, ai_analysis: str | None
) -> Slide:
    """Slide 2: Executive summary with stat cards and AI analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text_box(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        "الملخص التنفيذي",
        font_size=24,
        font_color=CLR_NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Stat cards — 4 across
    cards = [
        ("إجمالي العناوين", str(stats.total_records)),
        ("مباني قائمة", str(stats.total_with_buildings)),
        ("مباني مزالة", str(stats.total_without_buildings)),
        (
            "نسبة الإزالة",
            f"{stats.total_without_buildings / stats.total_records * 100:.0f}%",
        ),
    ]

    card_top = Inches(1.3)
    total_width = len(cards) * CARD_WIDTH + (len(cards) - 1) * CARD_SPACING
    start_left = (SLIDE_WIDTH - total_width) / 2

    for i, (label, value) in enumerate(cards):
        left = start_left + i * (CARD_WIDTH + CARD_SPACING)

        shape = slide.shapes.add_shape(1, left, card_top, CARD_WIDTH, CARD_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CLR_NAVY
        shape.line.fill.background()

        _add_text_box(
            slide,
            left + Inches(0.2),
            card_top + Inches(0.3),
            CARD_WIDTH - Inches(0.4),
            Inches(0.8),
            value,
            font_size=32,
            font_color=CLR_GOLD,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        _add_text_box(
            slide,
            left + Inches(0.2),
            card_top + Inches(1.2),
            CARD_WIDTH - Inches(0.4),
            Inches(0.5),
            label,
            font_size=14,
            font_color=CLR_WHITE,
            alignment=PP_ALIGN.CENTER,
        )

    if ai_analysis:
        _add_rich_text_box(
            slide,
            Inches(0.8),
            Inches(3.8),
            Inches(11.733),
            Inches(3.2),
            ai_analysis,
            base_font_size=13,
            font_color=CLR_DARK_GRAY,
        )

    return slide


def build_chart_slide(
    prs: Presentation,
    title: str,
    chart_buf: BytesIO,
    description: str | None = None,
) -> Slide:
    """Chart slide with optional AI description on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text_box(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        title,
        font_size=22,
        font_color=CLR_NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    if description:
        slide.shapes.add_picture(chart_buf, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT)

        desc_shape = slide.shapes.add_shape(1, DESC_LEFT, DESC_TOP, DESC_WIDTH, DESC_HEIGHT)
        desc_shape.fill.solid()
        desc_shape.fill.fore_color.rgb = CLR_LIGHT_GRAY
        desc_shape.line.fill.background()

        _add_text_box(
            slide,
            DESC_LEFT + Inches(0.3),
            DESC_TOP + Inches(0.3),
            DESC_WIDTH - Inches(0.6),
            Inches(0.5),
            "تحليل البيانات",
            font_size=14,
            font_color=CLR_GOLD,
            bold=True,
        )

        _add_rich_text_box(
            slide,
            DESC_LEFT + Inches(0.3),
            DESC_TOP + Inches(1.0),
            DESC_WIDTH - Inches(0.6),
            DESC_HEIGHT - Inches(1.3),
            description,
            base_font_size=12,
            font_color=CLR_DARK_GRAY,
        )
    else:
        slide.shapes.add_picture(
            chart_buf, FULL_CHART_LEFT, FULL_CHART_TOP, FULL_CHART_WIDTH, FULL_CHART_HEIGHT
        )

    return slide


def build_map_slide(
    prs: Presentation, map_buf: BytesIO, title: str = "خريطة مواقع المسح"
) -> Slide:
    """Full-width map slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text_box(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        title,
        font_size=22,
        font_color=CLR_NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    slide.shapes.add_picture(
        map_buf, FULL_CHART_LEFT, FULL_CHART_TOP, FULL_CHART_WIDTH, FULL_CHART_HEIGHT
    )

    return slide


def build_narrative_slide(prs: Presentation, title: str, text: str) -> Slide:
    """Text-only narrative slide (for insights, recommendations)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text_box(
        slide,
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
        title,
        font_size=24,
        font_color=CLR_NAVY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Gold accent line under title
    line = slide.shapes.add_shape(1, Inches(4), Inches(1.15), Inches(5.333), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_GOLD
    line.line.fill.background()

    # Content text (with markdown formatting)
    _add_rich_text_box(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(11.733),
        Inches(5.5),
        text,
        base_font_size=14,
        font_color=CLR_DARK_GRAY,
    )

    return slide
